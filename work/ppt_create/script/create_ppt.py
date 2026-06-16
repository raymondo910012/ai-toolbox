from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.shapes.placeholder import _InheritsDimensions

# === 常數 ===
TEMPLATE = r'/mnt/d/Kiro/project/ppt/Edgecore ppt format 2026_v0116_1_02.pptx'
FONT = 'Microsoft JhengHei'
CLR_TITLE = RGBColor(0x0F, 0x17, 0x2A)
CLR_BODY = RGBColor(0x0F, 0x17, 0x2A)
CLR_GREY = RGBColor(0x64, 0x74, 0x8B)
CLR_CARD_BG = RGBColor(0xF5, 0xF0, 0xFA)
CLR_CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)


def load_template():
    """載入 EC 模板並清除既有 slides"""
    prs = Presentation(TEMPLATE)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return prs


def add_cover(prs, title, subtitle=''):
    """新增封面/結尾頁（使用 cover layout，文字放下方）"""
    cover_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(cover_layout)
    box = s.shapes.add_textbox(Emu(914400), Emu(5000000), Emu(10332720), Emu(1000000))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CLR_TITLE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        box2 = s.shapes.add_textbox(Emu(914400), Emu(5700000), Emu(10332720), Emu(600000))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = FONT
        p2.font.size = Pt(24)
        p2.font.color.rgb = CLR_GREY
        p2.alignment = PP_ALIGN.CENTER
    return s


def add_content_slide(prs, title, body_lines):
    """新增內容頁（EC 模板背景 + Kiro 風格排版 + 淡紫卡片）"""
    content_layout = prs.slide_layouts[25]
    s = prs.slides.add_slide(content_layout)

    # 設定標題，移除不需要的 placeholder
    to_remove = []
    for shape in s.shapes:
        if isinstance(shape, _InheritsDimensions):
            idx = shape.placeholder_format.idx
            if idx == 0:
                tf = shape.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.name = FONT
                tf.paragraphs[0].font.size = Pt(30)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = CLR_TITLE
            elif idx in (2, 10):
                to_remove.append(shape)
    for shape in to_remove:
        shape._element.getparent().remove(shape._element)

    # 背景卡片（淡紫圓角矩形，send to back）
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(800000), Emu(1750000), Emu(10550000), Emu(4400000)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CLR_CARD_BG
    card.line.color.rgb = CLR_CARD_BORDER
    card.line.width = Pt(1)
    sp = card._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)

    # 內文
    box = s.shapes.add_textbox(Emu(914400), Emu(1828800), Emu(10332720), Emu(4500000))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(18)
        p.font.color.rgb = CLR_BODY
        p.space_after = Pt(10)
    return s


# === 使用範例 ===
if __name__ == '__main__':
    prs = load_template()

    add_cover(prs, '簡報標題', '副標題')

    add_content_slide(prs, 'Q1：題目標題', [
        '第一行內容',
        '第二行內容',
        '',
        '• bullet point'
    ])

    add_cover(prs, 'Thank You')

    prs.save('/mnt/d/Kiro/project/ppt/output.pptx')
    print('Done')
